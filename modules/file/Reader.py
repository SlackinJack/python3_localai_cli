# modules.file


import docx2txt as Docx2Txt
import json as JSON
import openpyxl as OpenPyXL
import pptx as PPTX
import PyPDF2 as PyPDF2
import subprocess as Subprocess


import modules.Configuration as Configuration
import modules.connection.response.AudioToText as AudioToText
import modules.connection.response.ImageToText as ImageToText
import modules.file.Operation as Operation
import modules.Print as Print
import modules.string.Path as Path
import modules.typecheck.TypeCheck as TypeCheck
import modules.typecheck.Types as Types
import modules.Util as Util


__extDOCX = []
__extPPTX = []
__extXLSX = []
__extPDF = []
__extAudio = []
__extImage = []


def loadConfiguration():
    global __extDOCX, __extPPTX, __extXLSX, __extPDF, __extAudio, __extImage
    readerConfig = Operation.readFile(Path.CONFIGS_READER_FILE_NAME, None, False)
    if readerConfig is not None:
        j = JSON.loads(readerConfig)
        __extDOCX = j.get("get_docx")
        __extPPTX = j.get("get_pptx")
        __extXLSX = j.get("get_xlsx")
        __extPDF = j.get("get_pdf")
        __extAudio = j.get("get_audio")
        __extImage = j.get("get_image")
    return


def __getFileMap():
    global __extDOCX, __extPPTX, __extXLSX, __extPDF, __extAudio, __extImage
    return {
        getDOCXText: __extDOCX,
        getPPTXText: __extPPTX,
        getXLSXText: __extXLSX,
        getPDFText: __extPDF,
        getAudioText: __extAudio,
        getImageText: __extImage
    }


def getDOCXText(filePath):
    TypeCheck.check(filePath, Types.STRING)
    return Docx2Txt.process(filePath)


def getPPTXText(filePath):
    TypeCheck.check(filePath, Types.STRING)
    content = ""
    for slide in PPTX.Presentation(filePath).slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if len(content) > 0:
                    content += " "
                content += shape.text
    return content


def getXLSXText(filePath):
    TypeCheck.check(filePath, Types.STRING)
    df = (OpenPyXL.load_workbook(filePath)).active
    rows = []
    for row in range(0, df.max_row):
        rowText = ""
        for col in df.iter_cols(1, df.max_column):
            if len(rowText) == 0:
                rowText += str(col[row].value)
            else:
                rowText += "," + str(col[row].value)
        rows.append(rowText)
    return Util.formatArrayToString(rows, "\n")


def getPDFText(filePath):
    TypeCheck.check(filePath, Types.STRING)
    pdfFile = PyPDF2.PdfReader(filePath)
    pdfPages = len(pdfFile.pages)
    pdfText = ""
    for i in range(pdfPages):
        pdfText += pdfFile.pages[i].extract_text()
    return pdfText


def getAudioText(filePath):
    TypeCheck.check(filePath, Types.STRING)
    return AudioToText.getAudioToTextResponse(filePath)


def getImageText(filePath):
    TypeCheck.check(filePath, Types.STRING)
    return ImageToText.getImageToTextResponse(Configuration.getConfig("image_to_text_prompt"), filePath)


def getFileExtension(filePath):
    TypeCheck.check(filePath, Types.STRING)
    if "." in filePath:
        f = filePath.split(".")
        return f[len(f) - 1]
    else:
        return ""


def getFileContents(filePath, writeIfNonexistent):
    TypeCheck.check(filePath, Types.STRING)
    TypeCheck.check(writeIfNonexistent, Types.BOOLEAN)
    fileExtension = getFileExtension(filePath)
    content = ""
    if len(fileExtension) > 0:
        for entry, value in __getFileMap().items():
            for ext in value:
                if fileExtension in ext:
                    functionCall = entry
                    content = functionCall(filePath)
                    if content is not None and len(content) > 0:
                        Util.printDump("\n" + content)
                        return content
                    else:
                        Util.printError("\nFile does not exist or content is empty.")
                        return None
    content = Operation.readFile(filePath, None, writeIfNonexistent)
    if content is not None:
        Util.printDump("\n" + content)
        return content
    else:
        Util.printError("\nFile does not exist or content is empty.")
        return None


# only supports linux
def openLocalFile(filePath, openerIn, shouldAsync):
    TypeCheck.check(filePath, Types.STRING)
    TypeCheck.checkList(openerIn, [Types.STRING, Types.NONE])
    TypeCheck.check(shouldAsync, Types.BOOLEAN)
    if filePath is not None and len(filePath) > 0:
        if openerIn is None:
            opener = ["xdg-open"]
        else:
            if " " in openerIn:
                opener = openerIn.split(" ")
            else:
                opener = [openerIn]
        if shouldAsync:
            Subprocess.Popen(opener + [filePath])
        else:
            Subprocess.call(opener + [filePath])
    return
