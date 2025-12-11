# package modules.core.file


import docx2txt as Docx2Txt
import json as JSON
import openpyxl as OpenPyXL
import pptx as PPTX
import PyPDF2 as PyPDF2
import subprocess as Subprocess


import modules.connection.response.AudioToText as AudioToText
import modules.connection.response.ImageToText as ImageToText
import modules.core.file.Operation as Operation
import modules.core.typecheck.TypeCheck as TypeCheck
import modules.core.typecheck.Types as Types
import modules.core.Util as Util
import modules.string.Path as Path
import modules.string.Prompt as Prompt


__fileMap = {}


def loadConfiguration():
    global __fileMap, __imageToTextUserPrompt
    readerConfig = Operation.readFile(Path.CONFIGS_READER_FILE_NAME, None, False)
    if readerConfig is not None:
        j = JSON.loads(readerConfig)
        for k in j.get("get_docx"):     __fileMap[k] = getDOCXText
        for k in j.get("get_pptx"):     __fileMap[k] = getPPTXText
        for k in j.get("get_xlsx"):     __fileMap[k] = getXLSXText
        for k in j.get("get_pdf"):      __fileMap[k] = getPDFText
        for k in j.get("get_audio"):    __fileMap[k] = getAudioText
        for k in j.get("get_image"):    __fileMap[k] = getImageText
    return


def __getFileMap():
    global __fileMap
    return __fileMap


def getDOCXText(filePath, promptIn=None):
    TypeCheck.enforce(filePath, Types.STRING)
    return Docx2Txt.process(filePath)


def getPPTXText(filePath, promptIn=None):
    TypeCheck.enforce(filePath, Types.STRING)
    content = ""
    for slide in PPTX.Presentation(filePath).slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if len(content) > 0:
                    content += " "
                content += shape.text
    return content


def getXLSXText(filePath, promptIn=None):
    TypeCheck.enforce(filePath, Types.STRING)
    df = (OpenPyXL.load_workbook(filePath)).active
    rows = []
    for row in range(0, df.max_row):
        rowText = ""
        for col in df.iter_cols(1, df.max_column):
            if len(rowText) == 0:
                rowText += str(col[row].value)
            else:
                rowText += "," + str(col[row].value)
        rows.append(rowText)
    return Util.formatArrayToString(rows, "\n")


def getPDFText(filePath, promptIn=None):
    TypeCheck.enforce(filePath, Types.STRING)
    pdfFile = PyPDF2.PdfReader(filePath)
    pdfPages = len(pdfFile.pages)
    pdfText = ""
    for i in range(pdfPages):
        pdfText += pdfFile.pages[i].extract_text()
    return pdfText


def getAudioText(filePath, promptIn=None):
    TypeCheck.enforce(filePath, Types.STRING)
    return AudioToText.getResponse(filePath)


def getImageText(filePath, promptIn=""):
    TypeCheck.enforce(filePath, Types.STRING)
    return ImageToText.getResponse(f"{Prompt.getImageToTextDefaultUserPrompt()} {promptIn}", filePath)


def getFileExtension(filePath):
    TypeCheck.enforce(filePath, Types.STRING)
    if "." in filePath:
        f = filePath.split(".")
        return f[len(f) - 1]
    else:
        return ""


def getFileContents(filePath, writeIfNonexistent, promptIn=None):
    TypeCheck.enforce(filePath, Types.STRING)
    TypeCheck.enforce(writeIfNonexistent, Types.BOOLEAN)
    fileExtension = getFileExtension(filePath)
    content = ""
    if len(fileExtension) > 0:
        opener = __getFileMap().get(fileExtension, None)
        if opener is not None:
            if promptIn is not None:
                content = opener(filePath, promptIn=promptIn)
            else:
                content = opener(filePath)
            if content is not None and len(content) > 0:
                Util.printDump(content)
                return content
            else:
                Util.printError("File does not exist or content is empty.")
                return None
    content = Operation.readFile(filePath, None, writeIfNonexistent)
    if content is not None:
        content = Util.cleanupString(content)
        Util.printDump(content)
        return content
    else:
        Util.printError("File does not exist or content is empty.")
        return None


# only supports linux
def openLocalFile(filePath, openerIn, shouldAsync):
    TypeCheck.enforce(filePath, Types.STRING)
    TypeCheck.enforceList(openerIn, [Types.STRING, Types.NONE])
    TypeCheck.enforce(shouldAsync, Types.BOOLEAN)
    if filePath is not None and len(filePath) > 0:
        if openerIn is None:
            opener = ["xdg-open"]
        else:
            if " " in openerIn:
                opener = openerIn.split(" ")
            else:
                opener = [openerIn]
        if shouldAsync:
            Subprocess.Popen(opener + [filePath])
        else:
            Subprocess.call(opener + [filePath])
    return
